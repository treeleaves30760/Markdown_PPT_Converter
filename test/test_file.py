from pptx import Presentation

presentation = Presentation()
for i in range(10):
    slide_layout = presentation.slide_layouts[i]
    slide = presentation.slides.add_slide(slide_layout)

presentation.save("test_presentation.pptx")
