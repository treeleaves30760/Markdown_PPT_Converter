import markdown
from pptx import Presentation
from pptx.util import Inches, Pt
from html.parser import HTMLParser

class MyHTMLParser(HTMLParser):
    def handle_data(self, data):
        self.data = data

class MarkdownToPptConverter:
    def __init__(self, md_file, ppt_file):
        self.md_file = md_file
        self.ppt_file = ppt_file

    def convert(self):
        with open(self.md_file, 'r') as f:
            md_content = f.read()

        html_content = markdown.markdown(md_content)

        # Initialize presentation
        presentation = Presentation()

        # Split html content by <h1> tag (assuming each slide starts with <h1>)
        slides_html = html_content.split('<h1>')

        for slide_html in slides_html:
            if slide_html.strip() == '':
                continue

            # Create a new slide
            slide_layout = presentation.slide_layouts[1]
            slide = presentation.slides.add_slide(slide_layout)

            # Get title and content
            title, content = slide_html.split('</h1>', 1)
            content = content.replace('<p>', '').replace('</p>', '\n')

            # Add title and content to slide
            title_placeholder = slide.shapes.title
            title_placeholder.text = title.strip()

            content_placeholder = slide.placeholders[1]
            content_placeholder.text = content.strip()

        # Save presentation
        presentation.save(self.ppt_file)

if __name__ == "__main__":
    converter = MarkdownToPptConverter('example.md', 'example.pptx')
    converter.convert()
