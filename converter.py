import markdown
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from html.parser import HTMLParser
import re
import requests
import io


class MyHTMLParser(HTMLParser):
    def handle_data(self, data):
        self.data = data


class MarkdownToPptConverter:
    def __init__(self, md_content, ppt_file, mode=0):
        """
        Initialize the converter with markdown content and PowerPoint file.
        The mode parameter specifies the type of input for the markdown content:
        0 - Markdown content is provided as a string
        1 - Markdown content is provided as a file path
        """
        self.md_content = md_content
        self.ppt_file = ppt_file
        self.mode = mode

    def download_image(self, url):
        """
        Download image from URL and return as BytesIO object.
        """
        response = requests.get(url)
        image = io.BytesIO(response.content)
        return image

    def convert(self):
        """
        Convert markdown content to PowerPoint presentation.
        The markdown content can be provided as a string or a file path.
        """
        # Initialize presentation
        presentation = Presentation()
        first_slide_created = False
        if self.mode == 1:
            with open(self.md_content_file, "r") as f:
                self.md_content = f.read()

            self.md_content = markdown.markdown(self.md_content)

        # Split markdown content by '---'
        slides_md = self.md_content.split("---")
        for slide_md in slides_md:
            slide_md = slide_md.strip()

            if slide_md == "":
                continue

            # Split title and content
            lines = slide_md.split("\n")
            title_line = lines[0].strip()
            content_lines = lines[1:]

            # Check if it's the first slide (cover page)
            if title_line.startswith("# ") and not first_slide_created:
                # Create cover slide
                slide_layout = presentation.slide_layouts[5]  # Title Slide
                slide = presentation.slides.add_slide(slide_layout)
                title = title_line.strip("# ").strip().replace("**", "")
                slide.shapes.title.text = title
                first_slide_created = True
            else:
                # Create a new slide for content
                slide_layout = presentation.slide_layouts[1]  # Title and Content
                slide = presentation.slides.add_slide(slide_layout)

                title = title_line.strip("# ").strip().replace("**", "")
                slide.shapes.title.text = title

                # Process content
                tf = slide.placeholders[1].text_frame
                tf.text = ""  # Clear default text

                for line in content_lines:
                    line = line.strip()
                    # Check Markdown or HTML image syntex
                    image_urls = re.findall(r"!\[.*?\]\((.*?)\)", line) + re.findall(
                        r'<img src="(.*?)"', line
                    )
                    for url in image_urls:
                        # Download image and add to slide
                        image_stream = self.download_image(url)
                        slide.shapes.add_picture(
                            image_stream, Inches(1), Inches(1), width=Inches(4)
                        )

                    line = line.replace("**", "")  # Remove '**' from line
                    if (
                        line.startswith("- ") and not image_urls
                    ):  # Check for bullet points
                        p = tf.add_paragraph()
                        p.text = line.strip("- ")
                        p.level = 0  # Adjust level as needed for nested bullets
                    elif (
                        line and not image_urls
                    ):  # Non-empty line that doesn't start with '-' and not an image
                        p = tf.add_paragraph()
                        p.text = line
                        p.level = 0

        # Save presentation
        presentation.save(self.ppt_file)


# Example usage of the MarkdownToPptConverter
if __name__ == "__main__":
    md_content = """ 
# Artificial Intelligence: An Overview

---

## Introduction

- Definition of AI
- Brief history of AI
- Importance of AI in the modern world

---

## Table of Contents

1. Core Features
2. Use Cases
3. Advantages
4. Limitations
5. Future Directions

---

## Core Features

- **Machine Learning**: Algorithms improve automatically through experience.
- **Natural Language Processing (NLP)**: Interaction between computers and humans using natural language.
- **Robotics**: Machines performing tasks with high precision and autonomy.
- **Computer Vision**: Enables machines to interpret and process visual data from the world.

---

## Use Cases

- **Healthcare**: Predictive analytics for patient care and disease diagnosis.
- **Finance**: Fraud detection and automated trading systems.
- **Manufacturing**: Predictive maintenance and optimized production lines.
- **Education**: Personalized learning experiences and automation of administrative tasks.

---

## Advantages

- **Efficiency and Speed**: Performing complex computations and data analyses rapidly.
- **Accuracy**: High precision in tasks like medical diagnoses and data entry.
- **Automation of Mundane Tasks**: Freeing up humans for creative and strategic roles.
- **Data Handling Capabilities**: Processing and interpreting vast amounts of data efficiently.

---

## Limitations

- **Ethical and Privacy Concerns**: Issues around data misuse and bias in AI algorithms.
- **High Costs**: Initial setup and maintenance of AI systems can be expensive.
- **Dependency and Job Displacement**: Potential for increased dependency on technology and displacement of jobs.
- **Complexity**: Designing, implementing, and maintaining AI systems require specialized knowledge.

---

## Future Directions

- **Ethical AI**: Developing guidelines and frameworks for the responsible use of AI.
- **Explainable AI**: Making AI decisions more transparent and understandable.
- **Integration into Everyday Life**: More seamless integration of AI in daily activities and industries.
- **Advancements in AI Technology**: Ongoing research to enhance AI's capabilities and reduce limitations.

---

"""
    converter = MarkdownToPptConverter(md_content, "example.pptx", mode=0)
    converter.convert()
    print("PowerPoint presentation created successfully!")
